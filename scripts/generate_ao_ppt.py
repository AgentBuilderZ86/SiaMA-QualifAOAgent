#!/usr/bin/env python3
from __future__ import annotations

import argparse
import io
import json
import re
import sys
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile
import xml.etree.ElementTree as ET

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError as exc:
    raise SystemExit("python-pptx est requis. Installez-le avec: python -m pip install python-pptx") from exc

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_CONFIG = ROOT / "config" / "ppt-sections.json"
DEFAULT_INPUT = ROOT / "data" / "ao-export.json"
DEFAULT_OUTPUT_DIR = ROOT / "output" / "ao-decks"
DEFAULT_ANALYSIS = ROOT / "output" / "template_analysis.json"
POTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main", "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
SIA_DARK = RGBColor(10, 21, 30)
SIA_BLUE = RGBColor(23, 48, 68)
SIA_TURQUOISE = RGBColor(0, 222, 204)
SIA_GREY = RGBColor(69, 86, 105)
SIA_LIGHT = RGBColor(244, 246, 252)
SIA_ORANGE = RGBColor(251, 174, 64)
SIA_FONT = "Sora-SIA"


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def text(value: Any, fallback: str = "À confirmer") -> str:
    value = "" if value is None else str(value).strip()
    return value or fallback


def as_list(value: Any, limit: int | None = None) -> list[str]:
    if isinstance(value, list):
        items = [text(item) for item in value if text(item, "")]
    elif value:
        items = [text(value)]
    else:
        items = []
    return items[:limit] if limit else items


def slug(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "_", value).strip("_")[:90] or "ao"


def load_config(path: Path) -> dict[str, Any]:
    if path.exists():
        return load_json(path)
    return {"defaultTemplatePath": str(Path.home() / "Documents" / "Modèles Office personnalisés" / "Copie de Sia_Template_Master.potx"), "profiles": {"standard": ["title", "agenda", "decision", "identification", "qualification", "signals", "risks", "watchpoints", "simulation", "proposal", "nextSteps", "sources", "conclusion"], "short": ["title", "decision", "identification", "signals", "watchpoints", "conclusion"], "finance": ["title", "decision", "identification", "simulation", "watchpoints", "sources", "conclusion"]}}


def normalized_template_source(template_path: Path) -> Path | io.BytesIO:
    if template_path.suffix.lower() != ".potx":
        return template_path
    buffer = io.BytesIO()
    with ZipFile(template_path, "r") as source, ZipFile(buffer, "w", ZIP_DEFLATED) as destination:
        for entry in source.infolist():
            content = source.read(entry.filename)
            if entry.filename == "[Content_Types].xml":
                content = content.replace(POTX_CONTENT_TYPE.encode("utf-8"), PPTX_CONTENT_TYPE.encode("utf-8"))
            destination.writestr(entry, content)
    buffer.seek(0)
    return buffer


def records_from_input(path: Path) -> list[dict[str, Any]]:
    raw = load_json(path)
    if isinstance(raw, list):
        return [item for item in raw if isinstance(item, dict)]
    if isinstance(raw, dict) and isinstance(raw.get("records"), list):
        return [item for item in raw["records"] if isinstance(item, dict)]
    return [raw] if isinstance(raw, dict) else []


def parse_ao_list(value: str | None) -> set[str]:
    if not value:
        return set()
    path = Path(value)
    raw = path.read_text(encoding="utf-8") if path.exists() else value
    return {item.strip().lower() for item in re.split(r"[\n,;]+", raw) if item.strip()}


def find_records(records: list[dict[str, Any]], ao_num: str | None, ao_list: str | None, all_records: bool) -> list[dict[str, Any]]:
    if all_records:
        return records
    expected = parse_ao_list(ao_list)
    if ao_num:
        expected.add(ao_num.strip().lower())
    if not expected:
        return records[:1]
    return [record for record in records if any(str(key or "").strip().lower() in expected for key in [record.get("ao", {}).get("aoNum"), record.get("ao", {}).get("displayAoNum"), record.get("aoNum")])]


def analyze_template(template_path: Path, output_path: Path) -> dict[str, Any]:
    fonts: Counter[str] = Counter()
    colors: Counter[str] = Counter()
    layouts = []
    with ZipFile(template_path) as archive:
        names = archive.namelist()
        presentation = ET.fromstring(archive.read("ppt/presentation.xml"))
        size = presentation.find(".//p:sldSz", NS)
        slide_size = {"cx": size.attrib.get("cx", ""), "cy": size.attrib.get("cy", "")} if size is not None else {"cx": "", "cy": ""}
        slide_count = len([name for name in names if re.match(r"ppt/slides/slide\d+\.xml$", name)])
        for name in sorted(item for item in names if re.match(r"ppt/slideLayouts/slideLayout\d+\.xml$", item)):
            root = ET.fromstring(archive.read(name))
            c_sld = root.find(".//p:cSld", NS)
            placeholders = [{"type": ph.attrib.get("type", ""), "idx": ph.attrib.get("idx", ""), "sz": ph.attrib.get("sz", "")} for ph in root.findall(".//p:ph", NS)]
            layouts.append({"file": name, "name": c_sld.attrib.get("name", "") if c_sld is not None else "", "placeholders": placeholders})
        for name in names:
            if not re.match(r"ppt/(slides|slideLayouts|theme)/.*\.xml$", name):
                continue
            root = ET.fromstring(archive.read(name))
            for latin in root.findall(".//a:latin", NS):
                if latin.attrib.get("typeface"):
                    fonts[latin.attrib["typeface"]] += 1
            for srgb in root.findall(".//a:srgbClr", NS):
                if srgb.attrib.get("val"):
                    colors[srgb.attrib["val"]] += 1
    analysis = {"templatePath": str(template_path), "generatedAt": datetime.now(timezone.utc).isoformat(), "slideSize": slide_size, "slideCount": slide_count, "layoutCount": len(layouts), "layouts": layouts, "dominantFonts": [{"name": name, "count": count} for name, count in fonts.most_common(12)], "dominantColors": [{"hex": name, "count": count} for name, count in colors.most_common(16)]}
    write_json(output_path, analysis)
    return analysis


def clear_template_slides(prs: Presentation) -> None:
    for slide_id in list(prs.slides._sldIdLst):  # noqa: SLF001
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)  # noqa: SLF001


def is_dark_layout(name: str) -> bool:
    normalized = name.lower()
    return any(token in normalized for token in ["navy", "dark", "black", "cover"])


def pick_layout(prs: Presentation, kind: str) -> Any:
    hints_by_kind = {
        "title": ["cover", "couverture", "titre", "title"],
        "agenda": ["agenda", "sommaire", "summary", "large-text", "text", "blank_title"],
        "text": ["large-text", "text", "texte", "content", "contenu", "blank_title", "titre"],
        "content": ["large-text", "text", "texte", "content", "contenu", "blank_title", "titre"],
        "table": ["table", "matrix", "matrice", "large-text", "content", "contenu", "blank_title"],
        "risks": ["table", "matrix", "matrice", "large-text", "content", "contenu", "blank_title"],
        "sources": ["table", "matrix", "matrice", "large-text", "content", "contenu", "blank_title"],
        "conclusion": ["large-text", "text", "conclusion", "closing", "titre", "content"],
    }
    hints = hints_by_kind.get(kind, hints_by_kind["text"])
    allow_dark = kind == "title"
    layouts = list(prs.slide_layouts)

    for hint in hints:
        for item in layouts:
            name = item.name.lower()
            if hint in name and (allow_dark or not is_dark_layout(name)):
                return item

    if not allow_dark:
        for item in layouts:
            if not is_dark_layout(item.name):
                return item

    return layouts[0]


def add_slide(prs: Presentation, kind: str = "content") -> Any:
    return prs.slides.add_slide(pick_layout(prs, kind))


def remove_shape(shape: Any) -> None:
    parent = shape._element.getparent()  # noqa: SLF001
    if parent is not None:
        parent.remove(shape._element)  # noqa: SLF001


def apply_text_style(paragraph: Any, size: int, bold: bool = False, color: RGBColor = SIA_DARK) -> None:
    paragraph.font.name = SIA_FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def set_shape_text(shape: Any, value: str | list[str], size: int = 13, bold: bool = False, color: RGBColor = SIA_DARK) -> None:
    frame = shape.text_frame
    frame.clear()
    values = value if isinstance(value, list) else [value]
    for index, item in enumerate(values or ["À confirmer"]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text(item)
        apply_text_style(paragraph, size, bold and index == 0, color)


def add_textbox(slide: Any, value: str, x: float, y: float, w: float, h: float, size: int = 13, bold: bool = False, color: RGBColor = SIA_DARK) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_text(shape, value, size, bold, color)
    return shape


def text_placeholders(slide: Any) -> list[Any]:
    return sorted([shape for shape in slide.shapes if getattr(shape, "is_placeholder", False) and getattr(shape, "has_text_frame", False)], key=lambda shape: (shape.top, shape.left))


def fill_placeholders(slide: Any, title: str, subtitle: str = "", bodies: list[str | list[str]] | None = None) -> None:
    text_shapes = text_placeholders(slide)
    used: set[int] = set()
    if text_shapes:
        set_shape_text(text_shapes[0], title, 22, True)
        used.add(id(text_shapes[0]))
    else:
        add_textbox(slide, title, 0.65, 0.42, 11.6, 0.45, 22, True)
    cursor = 1
    if subtitle and cursor < len(text_shapes):
        set_shape_text(text_shapes[cursor], subtitle, 10, False, SIA_GREY)
        used.add(id(text_shapes[cursor]))
        cursor += 1
    elif subtitle:
        add_textbox(slide, subtitle, 0.68, 0.92, 11.2, 0.35, 10, False, SIA_GREY)
    for body in bodies or []:
        if cursor < len(text_shapes):
            set_shape_text(text_shapes[cursor], body, 12)
            used.add(id(text_shapes[cursor]))
            cursor += 1
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False) and id(shape) not in used:
            remove_shape(shape)


def add_footer(slide: Any, payload: dict[str, Any]) -> None:
    sources = payload.get("sources", [])[:3]
    label = " | ".join(text(source.get("title")) for source in sources) if sources else "À confirmer"
    add_textbox(slide, f"Sources : {label}", 0.55, 7.12, 12.0, 0.22, 7, False, SIA_GREY)


def add_metric(slide: Any, label: str, value: str, x: float, y: float, color: RGBColor = SIA_TURQUOISE) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.55), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = SIA_LIGHT
    box.line.color.rgb = color
    add_textbox(slide, value, x + 0.15, y + 0.14, 2.25, 0.3, 17, True)
    add_textbox(slide, label, x + 0.15, y + 0.52, 2.25, 0.22, 8, False, SIA_GREY)


def add_bullets(slide: Any, bullets: list[str], x: float, y: float, w: float, h: float, size: int = 12) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets or ["À confirmer"]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text(bullet)
        apply_text_style(paragraph, size)


def add_table(slide: Any, rows: list[dict[str, Any]], x: float, y: float, w: float, h: float, max_rows: int = 8) -> None:
    if not rows:
        add_textbox(slide, "Données à confirmer", x, y, w, h, 13, False, SIA_GREY)
        return
    columns = list(rows[0].keys())[:5]
    visible = rows[:max_rows]
    table = slide.shapes.add_table(len(visible) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for col_index, column in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = text(column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SIA_DARK
        apply_text_style(cell.text_frame.paragraphs[0], 8, True, RGBColor(255, 255, 255))
    for row_index, row in enumerate(visible, start=1):
        for col_index, column in enumerate(columns):
            cell = table.cell(row_index, col_index)
            cell.text = text(row.get(column), "")
            apply_text_style(cell.text_frame.paragraphs[0], 8)


def can_chart(rows: list[dict[str, Any]]) -> bool:
    return bool(rows) and all(len([v for v in row.values() if isinstance(v, (int, float))]) == 1 and any(isinstance(v, str) and v.strip() for v in row.values()) for row in rows)


def add_chart_or_table(slide: Any, rows: list[dict[str, Any]], x: float, y: float, w: float, h: float) -> None:
    if can_chart(rows):
        data = CategoryChartData()
        data.categories = [text(next(value for value in row.values() if isinstance(value, str))) for row in rows]
        data.add_series("Valeur", [next(value for value in row.values() if isinstance(value, (int, float))) for row in rows])
        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    else:
        add_table(slide, rows, x, y, w, h)


def slide_title(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "title")
    ao = payload.get("ao", {})
    fill_placeholders(slide, text(ao.get("client")), text(ao.get("subject")), ["Qualification AO"])
    decision = payload.get("decision", {})
    add_metric(slide, "Recommandation", text(decision.get("recommendation")), 9.6, 1.25)
    add_metric(slide, "Score GO / NO GO", text(decision.get("score")), 9.6, 2.35, SIA_ORANGE)
    add_textbox(slide, f"AO {text(ao.get('displayAoNum'))} | {text(ao.get('deadline'))}", 0.65, 6.8, 8.0, 0.25, 9, False, SIA_GREY)


def slide_agenda(prs: Presentation, payload: dict[str, Any], sections: list[str]) -> None:
    labels = {"decision": "Décision GO / WATCH / NO GO", "identification": "Identification AO", "qualification": "Qualification synthétique", "signals": "Signaux de qualification", "risks": "Risques et questions ouvertes", "watchpoints": "Points de vigilance décision", "simulation": "Simulation financière", "proposal": "Storyboard propale", "nextSteps": "Prochaines étapes", "sources": "Sources et hypothèses", "conclusion": "Conclusion et prochaines actions"}
    slide = add_slide(prs, "agenda")
    fill_placeholders(slide, "Agenda", "Sections générées selon le profil de deck choisi", [[labels[item] for item in sections if item in labels]])
    add_footer(slide, payload)


def slide_identification(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "table")
    qual = payload.get("qualification", {})
    ident = qual.get("identification") or {}
    fill_placeholders(slide, "Identification AO", "Données opérationnelles utiles au comité de qualification")
    rows = [
        {"Champ": "Référence", "Valeur": text(ident.get("reference"))},
        {"Champ": "Maître d'ouvrage", "Valeur": text(ident.get("buyer"))},
        {"Champ": "Programme", "Valeur": text(ident.get("program"))},
        {"Champ": "Objet", "Valeur": text(ident.get("object"))},
        {"Champ": "Durée", "Valeur": text(ident.get("duration"))},
        {"Champ": "Date limite", "Valeur": text(ident.get("deadline"))},
        {"Champ": "Budget", "Valeur": text(ident.get("budget"))},
    ]
    add_table(slide, rows, 0.75, 1.55, 11.8, 4.6, 8)
    add_footer(slide, payload)


def slide_decision(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "text")
    decision = payload.get("decision", {})
    fill_placeholders(slide, "Décision IA", "Synthèse basée uniquement sur les données AO, la fiche qualification et les sources renseignées")
    add_metric(slide, "Recommandation", text(decision.get("recommendation")), 0.8, 1.65)
    add_metric(slide, "Score", text(decision.get("score")), 3.75, 1.65, SIA_ORANGE)
    add_metric(slide, "Confiance", text(decision.get("confidence")), 6.7, 1.65)
    add_bullets(slide, as_list(decision.get("reasons"), 4), 0.95, 3.05, 5.3, 2.6)
    add_bullets(slide, as_list(decision.get("vigilances"), 4), 6.9, 3.05, 5.3, 2.6)
    add_textbox(slide, f"Prochaine action : {text(decision.get('nextAction'))}", 0.95, 6.15, 11.1, 0.45, 12, True, SIA_BLUE)
    add_footer(slide, payload)


def slide_qualification(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "text")
    qual = payload.get("qualification", {})
    fill_placeholders(slide, "Qualification synthétique", text(qual.get("executiveSummary")))
    add_textbox(slide, "Contexte client", 0.75, 1.55, 5.6, 0.25, 12, True)
    add_textbox(slide, text(qual.get("clientContext")), 0.75, 1.9, 5.6, 1.1, 11, False, SIA_GREY)
    add_textbox(slide, "Périmètre", 0.75, 3.25, 5.6, 0.25, 12, True)
    add_textbox(slide, text(qual.get("scopeSynthesis")), 0.75, 3.6, 5.6, 1.1, 11, False, SIA_GREY)
    add_textbox(slide, "Enjeux et thèmes de gain", 6.85, 1.55, 5.5, 0.25, 12, True)
    add_bullets(slide, as_list(qual.get("businessIssues"), 4) + as_list(qual.get("winThemes"), 4), 6.85, 1.9, 5.5, 3.6)
    add_footer(slide, payload)


def slide_signals(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "table")
    qual = payload.get("qualification", {})
    fill_placeholders(slide, "Signaux de qualification", "Lecture GO / WATCH / NO GO des éléments détectés")
    rows = [
        {"Signal": text(item.get("label")), "Impact": text(item.get("impact")), "Détail": text(item.get("detail")), "Source": text(item.get("source"))}
        for item in qual.get("qualificationSignals", [])
        if isinstance(item, dict)
    ]
    add_table(slide, rows, 0.75, 1.55, 11.8, 4.6, 7)
    add_footer(slide, payload)


def slide_risks(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "risks")
    qual = payload.get("qualification", {})
    fill_placeholders(slide, "Risques et questions ouvertes", "Points à sécuriser avant engagement ferme")
    rows = [{"Risque": text(item.get("label")), "Sévérité": text(item.get("severity")), "Mitigation": text(item.get("mitigation"))} for item in qual.get("risks", []) if isinstance(item, dict)]
    add_table(slide, rows, 0.75, 1.55, 11.8, 2.55, 5)
    add_textbox(slide, "Questions de clarification", 0.75, 4.45, 5.6, 0.25, 12, True)
    add_bullets(slide, as_list(qual.get("clarificationQuestions"), 5), 0.75, 4.8, 11.4, 1.45, 11)
    add_footer(slide, payload)


def slide_watchpoints(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "table")
    qual = payload.get("qualification", {})
    fill_placeholders(slide, "Vigilances et décision manager", "Questions à trancher avant mobilisation")
    rows = [
        {"Point": text(item.get("point")), "Niveau": text(item.get("level")), "Question": text(item.get("question"))}
        for item in qual.get("decisionWatchpoints", [])
        if isinstance(item, dict)
    ]
    add_table(slide, rows, 0.75, 1.55, 11.8, 4.6, 7)
    manager = qual.get("managerRecommendation") or {}
    add_textbox(slide, f"Manager recommandé : {text(manager.get('primaryManager'))} — {text(manager.get('rationale'))}", 0.75, 6.25, 11.4, 0.45, 10, True, SIA_BLUE)
    add_footer(slide, payload)


def slide_simulation(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "table")
    sim = payload.get("simulation") or {}
    fill_placeholders(slide, "Simulation financière", text(sim.get("source"), "Données financières non renseignées"))
    add_metric(slide, "Total JH", text(sim.get("totalJours")), 0.8, 1.55)
    add_metric(slide, "Total HT", text(sim.get("totalHt")), 3.75, 1.55, SIA_ORANGE)
    add_metric(slide, "Total TTC", text(sim.get("totalTtc")), 6.7, 1.55)
    rows = sim.get("rows") if isinstance(sim.get("rows"), list) else []
    table_rows = [{"Phase": row.get("phase"), "Profil": row.get("profil"), "Jours": row.get("jours"), "TJM": row.get("tjm"), "Montant HT": row.get("montantHt")} for row in rows if isinstance(row, dict)]
    add_chart_or_table(slide, table_rows, 0.75, 2.95, 11.8, 3.3)
    add_footer(slide, payload)


def slide_proposal(prs: Presentation, payload: dict[str, Any]) -> None:
    proposal = payload.get("proposal") or {}
    storyboard = payload.get("storyboard") or []
    if proposal:
        slide = add_slide(prs, "text")
        fill_placeholders(slide, text(proposal.get("slideTitle"), "Section propale"), text(proposal.get("section"), "Contenu propale"))
        add_bullets(slide, as_list(proposal.get("keyMessages"), 5), 0.75, 1.6, 5.4, 2.4)
        add_textbox(slide, text(proposal.get("bodyText")), 6.75, 1.6, 5.55, 2.6, 12, False, SIA_GREY)
        add_footer(slide, payload)
    for item in storyboard[:2]:
        if isinstance(item, dict):
            slide = add_slide(prs, "text")
            fill_placeholders(slide, text(item.get("title")), text(item.get("keyMessage")))
            add_bullets(slide, as_list(item.get("bullets"), 6), 0.9, 1.65, 5.7, 4.4)
            add_textbox(slide, text(item.get("speakerNotes")), 7.0, 1.65, 5.3, 4.4, 12, False, SIA_GREY)
            add_footer(slide, payload)


def slide_next_steps(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "table")
    qual = payload.get("qualification", {})
    fill_placeholders(slide, "Prochaines étapes recommandées", "Plan d'action court pour passer de la qualification à la réponse")
    rows = [
        {"Action": text(item.get("action")), "Échéance": text(item.get("deadline")), "Owner": text(item.get("owner")), "Commande": text(item.get("workflowCommand"), "")}
        for item in qual.get("nextSteps", [])
        if isinstance(item, dict)
    ]
    add_table(slide, rows, 0.75, 1.55, 11.8, 4.6, 6)
    add_footer(slide, payload)


def slide_sources(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "sources")
    fill_placeholders(slide, "Sources et hypothèses", "Aucun chiffre n'est ajouté s'il n'est pas présent dans les données AO ou les référentiels")
    source_rows = [{"Source": text(item.get("title")), "URL": text(item.get("url"), "À confirmer"), "Consultation": text(item.get("consultedAt"))} for item in payload.get("sources", [])[:6] if isinstance(item, dict)]
    add_table(slide, source_rows, 0.75, 1.55, 11.8, 2.35, 6)
    qual = payload.get("qualification", {})
    add_textbox(slide, "Hypothèses à confirmer", 0.75, 4.25, 5.0, 0.25, 12, True)
    add_bullets(slide, as_list(qual.get("assumptions"), 5), 0.75, 4.6, 11.3, 1.5, 11)


def slide_conclusion(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = add_slide(prs, "conclusion")
    decision = payload.get("decision", {})
    qual = payload.get("qualification", {})
    fill_placeholders(slide, "Conclusion", f"Recommandation : {text(decision.get('recommendation'))}")
    add_bullets(slide, as_list(qual.get("differentiators"), 5), 0.95, 1.7, 5.4, 2.5, 14)
    add_textbox(slide, text(decision.get("nextAction")), 6.9, 1.7, 5.2, 2.0, 16, True, SIA_BLUE)
    add_footer(slide, payload)


SLIDE_BUILDERS = {"title": slide_title, "agenda": slide_agenda, "decision": slide_decision, "identification": slide_identification, "qualification": slide_qualification, "signals": slide_signals, "risks": slide_risks, "watchpoints": slide_watchpoints, "simulation": slide_simulation, "proposal": slide_proposal, "nextSteps": slide_next_steps, "sources": slide_sources, "conclusion": slide_conclusion}


def build_deck(payload: dict[str, Any], template_path: Path, output_path: Path, sections: list[str]) -> None:
    template_source = normalized_template_source(template_path)
    prs = Presentation(template_source if isinstance(template_source, io.BytesIO) else str(template_source))
    clear_template_slides(prs)
    for section in sections:
        builder = SLIDE_BUILDERS.get(section)
        if not builder:
            continue
        if section == "agenda":
            builder(prs, payload, sections)
        elif section == "simulation" and not payload.get("simulation"):
            continue
        else:
            builder(prs, payload)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def output_for(payload: dict[str, Any], output: Path, many: bool) -> Path:
    if output.suffix.lower() == ".pptx" and not many:
        return output
    ao = payload.get("ao", {})
    return output / f"{slug(text(ao.get('displayAoNum') or ao.get('aoNum')))}.pptx"


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Generer des decks AO a partir du template Sia.")
    parser.add_argument("--input", type=Path, default=DEFAULT_INPUT)
    parser.add_argument("--ao")
    parser.add_argument("--ao-list")
    parser.add_argument("--all", action="store_true")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT_DIR)
    parser.add_argument("--template", type=Path)
    parser.add_argument("--config", type=Path, default=DEFAULT_CONFIG)
    parser.add_argument("--profile", default="standard")
    parser.add_argument("--analyze-template", action="store_true")
    parser.add_argument("--analysis-output", type=Path, default=DEFAULT_ANALYSIS)
    args = parser.parse_args(argv)
    config = load_config(args.config)
    template_path = args.template or Path(config.get("defaultTemplatePath", ""))
    if not template_path.exists():
        raise SystemExit(f"Template introuvable: {template_path}")
    if args.analyze_template:
        analysis = analyze_template(template_path, args.analysis_output)
        print(f"Analyse template ecrite: {args.analysis_output} ({analysis['layoutCount']} layouts)")
        if not args.input.exists():
            return 0
    if not args.input.exists():
        raise SystemExit(f"Fichier d'entree introuvable: {args.input}")
    profiles = config.get("profiles", {})
    sections = profiles.get(args.profile) or profiles.get("standard") or ["title", "decision", "qualification", "conclusion"]
    records = find_records(records_from_input(args.input), args.ao, args.ao_list, args.all)
    if not records:
        raise SystemExit("Aucun AO correspondant aux criteres fournis.")
    generated = []
    for payload in records:
        target = output_for(payload, args.output, args.all or len(records) > 1)
        build_deck(payload, template_path, target, sections)
        generated.append(str(target))
    print(json.dumps({"generated": generated}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
